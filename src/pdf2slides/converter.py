"""PDF to PPTX Converter"""

import io
import math
import os
import re
from collections import defaultdict

import numpy as np
import pymupdf
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from sklearn.mixture import GaussianMixture
from sklearn.model_selection import GridSearchCV


class Converter:
    """Convert a PDF file to a PowerPoint file."""

    # Tolerance values in points (1 inch = 72 points)
    COLUMN_GAP_THRESHOLD = 13.2  # ~0.183" — min gap to detect separate columns
    SAME_BOX_GAP = 5.3           # ~0.074" — max vertical gap within same text box
    SAME_PARA_GAP = 2.6          # ~0.036" — max gap within same paragraph (wrapped text)
    MIN_INDENT = 0.35            # ~0.005" — min indent for sub-level detection

    def __init__(
        self,
        enable_ocr: bool = False,
        default_font: str | None = None,
        enforce_default_font: bool = False,
    ) -> None:
        """Initialize the converter.

        Args:
            enable_ocr: Whether to use OCR when extracting text.
            default_font: The font to use for all text.
            enforce_default_font: Whether to enforce the default font.
        """
        self.enable_ocr = enable_ocr
        self.default_font = default_font
        self.enforce_default_font = enforce_default_font
        self.ocr = None

        if self.enable_ocr:
            from paddleocr import PaddleOCR

            self.ocr = PaddleOCR(use_angle_cls=True, lang="en", show_log=False)

    @staticmethod
    def _validate_arguments(input_file_path: str, output_file_path: str) -> None:
        """Validate the input arguments."""

        if not os.path.exists(input_file_path):
            raise FileNotFoundError(
                "Input file '{}' does not exist".format(input_file_path)
            )

        output_directory = os.path.dirname(output_file_path)
        if output_directory and not os.path.exists(output_directory):
            os.makedirs(output_directory)

    def convert(self, input_file_path: str, output_file_path: str) -> None:
        """Convert a PDF to a PowerPoint.

        Args:
            input_file_path: The path of the input `.pdf` file.
            output_file_path: The path of the generated `.pptx` file.
        """

        # Validate the arguments. Raise an error if they are invalid.
        Converter._validate_arguments(input_file_path, output_file_path)

        pptx_output = Presentation()

        with pymupdf.open(input_file_path) as pdf_document:
            # Determine slide size from the first page
            first_page = pdf_document[0]
            first_page_width = first_page.rect.width
            first_page_height = first_page.rect.height

            # Convert PDF page size from points to inches
            slide_width = first_page_width / 72.0
            slide_height = first_page_height / 72.0

            # Set the slide size as the matching size
            pptx_output.slide_width = Inches(slide_width)
            pptx_output.slide_height = Inches(slide_height)

            all_contents = [
                self._get_page_contents(pdf_page) for pdf_page in pdf_document
            ]

            # No page contains editable text. The document is likely scanned.
            scanned_document = self.enable_ocr and all(
                (not page_contents["text_blocks"]) for page_contents in all_contents
            )

            self._construct_pptx(
                pdf_document, all_contents, pptx_output, scanned_document
            )

        output_directory = os.path.dirname(output_file_path)
        if output_directory and not os.path.exists(output_directory):
            os.makedirs(output_directory)

        # Save PowerPoint presentation
        pptx_output.save(output_file_path)

    def _get_page_contents(self, pdf_page: pymupdf.Page) -> dict:
        """Get all editable contents (excluding OCR'ed text) from page."""

        page_contents = {}
        page_contents["text_blocks"] = self._get_page_text(pdf_page)
        page_contents["drawings"] = self._get_page_drawings(pdf_page)
        page_contents["images"] = self._get_page_images(pdf_page)
        page_contents["xref_smask_map"] = self._get_page_xref_smask_map(pdf_page)
        return page_contents

    def _get_page_text(self, pdf_page: pymupdf.Page) -> list[dict]:
        """Get all text blocks from the page."""

        blocks = pdf_page.get_textpage().extractDICT()["blocks"]
        text_blocks = [block for block in blocks if block["type"] == 0]
        return text_blocks

    def _get_page_drawings(self, pdf_page: pymupdf.Page) -> list[dict]:
        """Get all drawings from the page."""

        drawings = pdf_page.get_drawings()
        return drawings

    def _get_page_images(self, pdf_page: pymupdf.Page) -> list[dict]:
        """Get all images from the page."""

        images = pdf_page.get_image_info(xrefs=True)
        return images

    def _get_page_xref_smask_map(self, pdf_page: pymupdf.Page) -> dict[int, int]:
        """Store the smask of the images in a dictionary."""

        xref_smask_map = {item[0]: item[1] for item in pdf_page.get_images()}
        return xref_smask_map

    def _get_page_ocr_text(self, pdf_page: pymupdf.Page) -> list:
        """OCR the page and get the result."""

        page_pixmap = pdf_page.get_pixmap(dpi=300)
        page_bytes = page_pixmap.tobytes()
        ocr_result = self.ocr.ocr(page_bytes, cls=False)
        return ocr_result

    @staticmethod
    def _readable_font_name(raw_name: str, fallback: str) -> str:
        """Clean up PyMuPDF font names like 'Type3 (179 0 R)' or 'ABCDEF+Calibri'."""
        cleaned = raw_name
        # Remove embedded object references: "Type3 (179 0 R)" → "Type3"
        cleaned = re.sub(r'\s*\(\d+\s+\d+\s+R\)\s*$', '', cleaned)
        # Remove prefix hex garbage: "ABCDEF+Calibri" → "Calibri"
        cleaned = re.sub(r'^[A-F0-9]{6}\+', '', cleaned)
        # Generic type names get the fallback
        if cleaned.lower().startswith('type'):
            return fallback
        return cleaned

    @staticmethod
    def _detect_alignment(text_block: dict) -> int:
        """Detect text alignment from span positions within the block."""
        block_x0, _, block_x1, _ = text_block["bbox"]
        block_width = block_x1 - block_x0
        if block_width <= 0:
            return PP_ALIGN.LEFT

        total_margin = 0.0
        count = 0
        for line in text_block["lines"]:
            for span in line["spans"]:
                sx0 = span["bbox"][0]
                # Normalised left margin: 0 = block left, 1 = block right
                margin = (sx0 - block_x0) / block_width
                total_margin += margin
                count += 1

        if count == 0:
            return PP_ALIGN.LEFT

        avg_margin = total_margin / count
        # If spans are significantly inset from the left edge → right-aligned
        if avg_margin > 0.3:
            return PP_ALIGN.RIGHT
        return PP_ALIGN.LEFT

    # ------------------------------------------------------------------
    # Block merging (to fix pymupdf splitting paragraphs into separate blocks)
    # ------------------------------------------------------------------

    def _merge_text_blocks(self, text_blocks: list[dict]) -> list[dict]:
        """Merge adjacent blocks that belong together into one text box.

        Strategy:
        1. Cluster blocks into columns by x-position
        2. Within each column, merge vertically adjacent blocks (small gap)
        3. Return merged blocks
        """
        if not text_blocks:
            return []

        # Sort by y-position (top-to-bottom)
        sorted_blocks = sorted(text_blocks, key=lambda b: (b["bbox"][1], b["bbox"][0]))

        # Cluster into columns by x-position
        columns = self._cluster_columns(sorted_blocks)

        # Merge adjacent blocks within each column
        merged = []
        for col in columns:
            col_sorted = sorted(col, key=lambda b: b["bbox"][1])
            groups = self._group_adjacent(col_sorted)
            for group in groups:
                merged_block = self._combine_blocks(group)
                if merged_block:
                    merged.append(merged_block)

        return merged

    def _cluster_columns(self, blocks: list[dict]) -> list[list[dict]]:
        """Group blocks into columns by x-position clustering."""
        if not blocks:
            return []

        # Get center x of each block
        block_xs = []
        for b in blocks:
            x0, _, x1, _ = b["bbox"]
            cx = (x0 + x1) / 2.0
            block_xs.append(cx)

        # Sort unique x positions
        unique_xs = sorted(set(round(cx / 3.0) * 3.0 for cx in block_xs))

        # Cluster
        clusters = []
        current = [unique_xs[0]]
        for x in unique_xs[1:]:
            if x - current[-1] < self.COLUMN_GAP_THRESHOLD:
                current.append(x)
            else:
                clusters.append((min(current), max(current)))
                current = [x]
        clusters.append((min(current), max(current)))

        # Assign each block to the nearest cluster
        result = [[] for _ in clusters]
        for block in blocks:
            x0, _, x1, _ = block["bbox"]
            cx = (x0 + x1) / 2.0
            best_idx = 0
            best_dist = float("inf")
            for i, (cmin, cmax) in enumerate(clusters):
                cc = (cmin + cmax) / 2.0
                d = abs(cx - cc)
                if d < best_dist:
                    best_dist = d
                    best_idx = i
            result[best_idx].append(block)

        return result

    def _group_adjacent(self, blocks: list[dict]) -> list[list[dict]]:
        """Group vertically adjacent blocks within a column."""
        if not blocks:
            return []

        groups = [[blocks[0]]]
        for block in blocks[1:]:
            prev_bottom = groups[-1][-1]["bbox"][3]
            this_top = block["bbox"][1]
            gap = this_top - prev_bottom
            if gap < self.SAME_BOX_GAP:
                groups[-1].append(block)
            else:
                groups.append([block])
        return groups

    def _combine_blocks(self, blocks: list[dict]) -> dict | None:
        """Combine a list of blocks into one merged block with all lines."""
        if not blocks:
            return None

        if len(blocks) == 1:
            return blocks[0]

        # Compute combined bbox
        x0 = min(b["bbox"][0] for b in blocks)
        y0 = min(b["bbox"][1] for b in blocks)
        x1 = max(b["bbox"][2] for b in blocks)
        y1 = max(b["bbox"][3] for b in blocks)

        # Collect all lines in order
        all_lines = []
        for block in blocks:
            all_lines.extend(block["lines"])

        merged = {
            "type": 0,
            "bbox": (x0, y0, x1, y1),
            "lines": all_lines,
        }
        return merged

    # ------------------------------------------------------------------
    # Text box creation
    # ------------------------------------------------------------------

    def _add_text_block_to_slide(self, text_block: dict, slide) -> None:
        """Add all text in the block to a single text box.

        One box per block, one paragraph per line, one run per span.
        """
        if not text_block.get("lines"):
            return

        # Block bounding box (PDF points)
        bx0, by0, bx1, by1 = text_block["bbox"]
        if bx1 - bx0 < 1 or by1 - by0 < 1:
            return

        # Calculate position on the slide (points → inches)
        text_box_left = Inches(bx0 / 72.0)
        text_box_top = Inches(by0 / 72.0)
        text_box_width = Inches((bx1 - bx0) / 72.0)
        text_box_height = Inches((by1 - by0) / 72.0)

        # Add single text box for the whole block
        new_text_box = slide.shapes.add_textbox(
            text_box_left, text_box_top, text_box_width, text_box_height
        )

        text_frame = new_text_box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.margin_bottom = 0
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        text_frame.word_wrap = True

        # Determine alignment from block geometry (right-aligned page numbers, etc.)
        alignment = self._detect_alignment(text_block)

        lines = text_block["lines"]
        fallback_font = (
            self.default_font if self.default_font else "Calibri"
        )

        for li, line in enumerate(lines):
            if li == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.alignment = alignment

            # Infer spacing from PDF line positions
            if li > 0:
                prev_line = lines[li - 1]
                prev_bottom = prev_line["bbox"][3]
                this_top = line["bbox"][1]
                line_gap = this_top - prev_bottom
                # Calculate max font size on this line for relative spacing
                max_size = max(
                    (s["size"] for s in line["spans"] if s["size"] >= 1),
                    default=12,
                )
                if max_size > 0 and line_gap > 0:
                    # Set space_before to replicate PDF line spacing
                    p.space_before = Pt(line_gap)

            # Add spans as runs within the paragraph
            for si, span in enumerate(line["spans"]):
                span_text = span["text"]
                span_font_size = span["size"]

                if not span_text or span_font_size < 1:
                    continue

                span_font_name = span["font"]
                span_font_color = span["color"]
                span_font_is_italic = bool(span["flags"] & 2**1)
                span_font_is_bold = bool(span["flags"] & 2**4)

                # Reuse existing run for first span, add new ones after
                if si == 0 and p.runs:
                    run = p.runs[0]
                else:
                    run = p.add_run()

                run.text = span_text
                run.font.size = Pt(span_font_size)

                # Apply font name (cleaned from Type3 / hex prefixes)
                run.font.name = (
                    self.default_font
                    if self.enforce_default_font and self.default_font
                    else self._readable_font_name(span_font_name, fallback_font)
                )

                # Apply colour (skip gradient cases that turn black)
                try:
                    run.font.color.rgb = RGBColor.from_string(
                        f"{span_font_color:06X}"
                    )
                except (ValueError, AttributeError):
                    pass

                run.font.italic = span_font_is_italic
                run.font.bold = span_font_is_bold

    def _add_image_to_slide(
        self,
        pdf_doc: pymupdf.Document,
        image: dict,
        smask: int,
        bbox: pymupdf.Rect,
        slide,
    ) -> None:
        """Add the image to a slide."""

        xref = image["xref"]
        base_image = pymupdf.Pixmap(pdf_doc, xref)

        # The image itself is a smask. We ingore this image.
        if not base_image.colorspace:
            return

        # The transparency data is contained in the smask
        if smask:
            smask_pixmap = pymupdf.Pixmap(pdf_doc, smask)
            # The alpha channel is already present but contains no information
            if base_image.alpha:
                base_image.set_alpha(smask_pixmap.samples)
            # The alpha channel is not present
            else:
                base_image = pymupdf.Pixmap(base_image, smask_pixmap)

        transformation_matrix = pymupdf.Matrix(image["transform"])
        transformed_image_bytes = Converter._transform_image(
            base_image, bbox, transformation_matrix
        )

        # Calculate position on the slide based on original PDF coordinates
        left = Inches(bbox[0] / 72.0)
        top = Inches(bbox[1] / 72.0)
        width = Inches((bbox[2] - bbox[0]) / 72.0)

        # Add image to the PowerPoint slide
        slide.shapes.add_picture(
            io.BytesIO(transformed_image_bytes), left, top, width=width
        )

    def _add_drawing_to_slide(
        self,
        drawing: dict,
        bbox: pymupdf.Rect,
        pdf_page_width: float,
        pdf_page_height: float,
        slide,
    ) -> None:
        """Add all drawings to a slide.

        Reference: https://pymupdf.readthedocs.io/en/latest/recipes-drawing-and-graphics.html
        """
        # BUG: Drawings with gradient color will be lost

        # Create a temporary page with the same dimensions to serve as a drawing canvas for the shape
        # NOTE: Empirically, creating a new page for each shape is no slower than drawing each shape and then removing it on the same page
        temp_pdf = pymupdf.open()
        temp_page = temp_pdf.new_page(width=pdf_page_width, height=pdf_page_height)

        shape = temp_page.new_shape()
        for item in drawing["items"]:  # these are the draw commands
            if item[0] == "l":  # line
                shape.draw_line(item[1], item[2])
            elif item[0] == "re":  # rectangle
                shape.draw_rect(item[1])
            elif item[0] == "qu":  # quad
                shape.draw_quad(item[1])
            elif item[0] == "c":  # curve
                shape.draw_bezier(item[1], item[2], item[3], item[4])
            elif item[0] == "cu":  # curve
                shape.draw_curve(item[1], item[2], item[3])
            elif item[0] == "o":  # circle/ellipse
                if item[1] == 1:
                    shape.draw_circle(item[2], item[3])
                else:
                    shape.draw_oval(item[2], item[3])
            elif item[0] == "s":  # section of circle/elipse
                shape.draw_sector(item[1], item[2], item[3])
        shape.finish(
            fill=drawing.get("fill"),
            color=drawing.get("color"),
            width=drawing.get("width"),
            dashes=drawing.get("dashes"),
            even_odd=drawing.get("even_odd", True),
            closePath=drawing.get("closePath", True),
            fill_opacity=drawing.get("fill_opacity", 1),
            stroke_opacity=drawing.get("stroke_opacity", 1),
        )

        image_bytes = temp_page.get_pixmap().tobytes("png")

        # Calculate position on the slide based on original PDF coordinates
        left = Inches(bbox[0] / 72.0)
        top = Inches(bbox[1] / 72.0)
        width = Inches((bbox[2] - bbox[0]) / 72.0)

        try:
            # Add drawing to the PowerPoint slide
            slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, width=width)
        except:
            # Something went wrong. Ignore this drawing.
            pass

    def _add_ocr_to_slide(self, ocr_line: list, fontsize: float, slide) -> None:
        """Add all text obtained from ocr to a slide."""

        # Since the input pixmap has dpi=300 rather than dpi=72 by default, we need to set the scaling factor as 72/300 = 0.24
        scaling_factor = 0.24

        text = ocr_line[1][0]
        rect = ocr_line[0]

        # Get bounding box coordinates
        left, top, right, bottom = Converter._get_bbox(rect)
        left *= scaling_factor
        top *= scaling_factor
        right *= scaling_factor
        bottom *= scaling_factor

        # Calculate position and size in PowerPoint slide
        x, y = left, top
        width, height = right - left, bottom - top

        # Add text box to slide
        text_box = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(width), Pt(height))
        text_frame = text_box.text_frame
        text_frame.text = text

        # Set the margin to 0 to visually mimic original layout
        text_frame.margin_bottom = 0
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0

        for paragraph in text_frame.paragraphs:
            # Set font size
            paragraph.font.size = Pt(fontsize)
            # Set the font family to the default font
            if self.default_font:
                paragraph.font.name = self.default_font

    def _construct_pptx(
        self,
        pdf_document: pymupdf.Document,
        all_contents: list[dict],
        pptx_output: Presentation,
        scanned_document: bool,
    ) -> None:
        """Construct output slides with scanned PDF file."""

        if scanned_document:
            ocr_results = [
                self._get_page_ocr_text(pdf_page) for pdf_page in pdf_document
            ]
            detected_sizes = [
                Converter._get_ocr_detected_font_size(line)
                for ocr_result in ocr_results
                for line in ocr_result[0]
            ]
            suitable_sizes = Converter._get_suitable_sizes(detected_sizes)

        for page_num, pdf_page in enumerate(pdf_document):
            pdf_page_width = pdf_page.rect.width
            pdf_page_height = pdf_page.rect.height

            # Create a new slide with blank layout
            slide_layout = pptx_output.slide_layouts[6]  # Blank layout
            slide = pptx_output.slides.add_slide(slide_layout)

            page_content = all_contents[page_num]
            drawings: list[dict] = page_content["drawings"]
            drawing_bboxes: list[pymupdf.Rect] = [
                drawing["rect"] for drawing in drawings
            ]
            images: list[dict] = page_content["images"]
            image_bboxes = [pymupdf.Rect(image["bbox"]) for image in images]
            xref_smask_map = page_content["xref_smask_map"]
            all_shapes = drawings + images
            all_bboxes = drawing_bboxes + image_bboxes

            # Sort indices of shapes by area in descending order
            indices = Converter._sort_shapes_by_area(all_bboxes)

            # Add the shapes to the slide. Smaller shapes will appear in front.
            for index in indices:
                shape = all_shapes[index]
                bbox = all_bboxes[index]
                # The shape is a drawing
                if index < len(drawings):
                    # Do not add svg characters to the slides
                    if scanned_document and Converter._drawing_is_character(shape):
                        continue
                    self._add_drawing_to_slide(
                        shape, bbox, pdf_page_width, pdf_page_height, slide
                    )
                # The shape is an image
                else:
                    xref = shape["xref"]
                    # This is an inline image
                    if xref == 0:
                        continue
                    # Do not add pure text images to slide
                    if scanned_document and self._image_is_pure_text(
                        pdf_document, xref
                    ):
                        continue
                    smask = xref_smask_map[xref]
                    self._add_image_to_slide(pdf_document, shape, smask, bbox, slide)

            # Iterate through text blocks and add them to the slide
            # NOTE: We add text blocks to slides AFTER we add images in order that they appear in front of the images.
            text_blocks: list[dict] = page_content["text_blocks"]

            # NEW: Merge adjacent blocks before creating text boxes
            merged_blocks = self._merge_text_blocks(text_blocks)

            for merged_block in merged_blocks:
                self._add_text_block_to_slide(merged_block, slide)

            if scanned_document:
                page_ocr_results = ocr_results[page_num]
                # NOTE: The index 0 is necessary
                for line_num, ocr_line in enumerate(page_ocr_results[0]):
                    fontsize = suitable_sizes[
                        sum(len(page_ocr[0]) for page_ocr in ocr_results[:page_num])
                        + line_num
                    ]
                    self._add_ocr_to_slide(ocr_line, fontsize, slide)

    @staticmethod
    def _sort_shapes_by_area(all_bboxes: list[pymupdf.Rect]) -> list[int]:
        """Sort indices of shapes by area in descending order."""
        indices = sorted(
            range(len(all_bboxes)),
            key=lambda i: all_bboxes[i].get_area(),
            reverse=True,
        )
        return indices

    @staticmethod
    def _get_bbox(rect: list[list[float]]) -> list[float]:
        tl, tr, br, bl = rect
        left = min(tl[0], bl[0])
        top = min(tl[1], tr[1])
        right = max(tr[0], br[0])
        bottom = max(br[1], tr[1])
        return left, top, right, bottom

    @staticmethod
    def _get_ocr_detected_font_size(ocr_line: list) -> float:
        """Calculate the detected font size as the minimum of bbox height and width, clamped between 1 and 64."""

        # Since the input pixmap has dpi=300 rather than dpi=72 by default, we need to set the scaling factor as 72/300 = 0.24
        scaling_factor = 0.24

        rect = ocr_line[0]
        # Get bounding box coordinates
        left, top, right, bottom = Converter._get_bbox(rect)
        left *= scaling_factor
        top *= scaling_factor
        right *= scaling_factor
        bottom *= scaling_factor

        height = bottom - top
        width = right - left

        # Clamp the fontsize between 1 and 64 to avoid extreme fontsizes
        fontsize = np.clip(min(height, width), 1, 64)
        return fontsize

    @staticmethod
    def _gmm_bic_score(estimator, X):
        """Callable to pass to GridSearchCV that will use the BIC score.

        Reference: https://scikit-learn.org/stable/auto_examples/mixture/plot_gmm_selection.html#sphx-glr-auto-examples-mixture-plot-gmm-selection-py
        """
        # Make it negative since GridSearchCV expects a score to maximize
        return -estimator.bic(X)

    @staticmethod
    def _get_suitable_sizes(detected_sizes: list[float]) -> list[float]:
        """Determine the suitable font sizes using GMM clustering.

        Reference: https://scikit-learn.org/stable/auto_examples/mixture/plot_gmm_selection.html#sphx-glr-auto-examples-mixture-plot-gmm-selection-py
        """

        # Keep the detected sizes as they are, otherwise an error will be raised
        if len(detected_sizes) < 5:
            return detected_sizes

        # Convert the list into a numpy column vector
        X = np.asarray(detected_sizes).reshape((-1, 1))

        # Use grid search to determine the optimal number of components
        param_grid = {
            "n_components": range(1, 7),
        }
        grid_search = GridSearchCV(
            GaussianMixture(init_params="k-means++"),
            param_grid=param_grid,
            scoring=Converter._gmm_bic_score,
        )
        grid_search.fit(X)

        # Calculate the suitable sizes as the predicted cluster means
        best_estimator = grid_search.best_estimator_
        labels = best_estimator.predict(X)
        suitable_sizes = list(map(lambda i: best_estimator.means_[i], labels))
        return suitable_sizes

    @staticmethod
    def _transform_image(
        base_image: pymupdf.Pixmap,
        bbox: pymupdf.Rect,
        transformation_matrix: pymupdf.Matrix,
    ) -> bytes:
        """Transform the image according to the transformation matrix and return the transformed image in bytes."""

        # Convert the colorspace to RGB explicitly
        base_image = pymupdf.Pixmap(pymupdf.Colorspace(pymupdf.CS_RGB), base_image)

        # Add a fully opaque alpha channel if the channel is not present
        if not base_image.alpha:
            alpha = bytearray(255 for _ in range(base_image.w * base_image.h))
            base_image = pymupdf.Pixmap(base_image, 1)
            base_image.set_alpha(alpha)

        # Convert the base image to a Pillow Image object
        image_data = base_image.samples
        pil_image = Image.frombytes("RGBA", (base_image.w, base_image.h), image_data)

        # Define a scaling factor to enhance the output quality
        scaling_factor = 4

        # Calculate the matrix that transforms the original image to the displayed image
        shrink = pymupdf.Matrix(
            scaling_factor / base_image.w, 0, 0, scaling_factor / base_image.h, 0, 0
        )
        translation = pymupdf.Matrix(1, 0, 0, 1, -bbox[0], -bbox[1])
        transformation_matrix = shrink * transformation_matrix * translation
        transformation_matrix.invert()

        # Obtain the transformed image
        pil_image = pil_image.transform(
            (
                scaling_factor * int(bbox[2] - bbox[0]),
                scaling_factor * int(bbox[3] - bbox[1]),
            ),
            Image.AFFINE,
            # The parameters in PIL have a different order than in PyMuPDF
            (
                transformation_matrix.a,
                transformation_matrix.c,
                transformation_matrix.e,
                transformation_matrix.b,
                transformation_matrix.d,
                transformation_matrix.f,
            ),
            resample=Image.Resampling.BICUBIC,
        )

        # Convert Pillow image back to bytes
        output_bytes = io.BytesIO()
        pil_image.save(output_bytes, format="PNG")
        return output_bytes.getvalue()

    @staticmethod
    def _drawing_is_character(drawing: dict) -> bool:
        """Determine whether the drawing is a character using heuristics."""

        command_length_threshold = 6
        return all(
            [
                not drawing["closePath"],
                drawing["type"] == "f",
                len(drawing["items"]) > command_length_threshold,
            ]
        )

    def _image_is_pure_text(self, pdf_doc: pymupdf.Document, image_xref: dict) -> bool:
        """Use the layout engine to determine whether the image is pure text."""

        base_image = pymupdf.Pixmap(pdf_doc, image_xref)
        image_bytes = base_image.tobytes()
        layout_result = self.layout_engine(image_bytes)

        # The image is likely decorative. It does not contain text.
        if not layout_result:
            return False

        # Any area type other than figure is considered as text
        for area in layout_result:
            if area["type"] == "Figure":
                return False
        return True

    @staticmethod
    def _validate_arguments(input_file_path: str, output_file_path: str) -> None:
        if not os.path.isfile(input_file_path):
            raise FileNotFoundError(f"The file {input_file_path} does not exist.")

        input_file_extension = os.path.splitext(input_file_path)[1]
        if input_file_extension != ".pdf":
            raise ValueError(
                f"Invalid input file extension. Expected a .pdf file, but got {input_file_extension}"
            )

        output_file_extension = os.path.splitext(output_file_path)[1]
        if output_file_extension != ".pptx":
            raise ValueError(
                f"Invalid output file extension. Expected a .pptx file, but got {output_file_extension}"
            )
